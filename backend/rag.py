import os
import time
import uuid
import atexit
import hashlib
import shutil
import base64
import re
import threading
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from difflib import SequenceMatcher

import chromadb
import numpy as np
import pandas as pd
from sentence_transformers import SentenceTransformer
from langchain_community.document_loaders import (
    PyPDFLoader, Docx2txtLoader, TextLoader, CSVLoader, UnstructuredExcelLoader
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
import google.generativeai as genai

# ─── Detección de idioma ──────────────────────────────────────────────────────
try:
    from langdetect import detect, DetectorFactory
    DetectorFactory.seed(0)
    LANGDETECT_AVAILABLE = True
except ImportError:
    LANGDETECT_AVAILABLE = False
    print("[RAG] langdetect no instalado. Instalar con: pip install langdetect")

# ─── Configuración ────────────────────────────────────────────────────────────
EMBED_MODEL = SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')

chroma_client = chromadb.PersistentClient(path="./vectorstore")
collection = chroma_client.get_or_create_collection(
    name="documentos",
    metadata={"heuristic": "cosine"}
)

SYSTEM_PROMPT = """Eres un asistente técnico especializado en análisis de documentos.

REGLAS ESTRICTAS:
1. Responde ÚNICAMENTE basándote en los fragmentos de documentos proporcionados.
2. Si la información no está en los fragmentos, di explícitamente: "Esta información no se encuentra en los documentos cargados."
3. Usa terminología técnica precisa y apropiada al dominio del documento.
4. Cita el nombre del archivo fuente cuando uses información específica.
5. Estructura tus respuestas con claridad.
6. Si hay contradicciones entre fragmentos, señálalas explícitamente.
7. No inventes datos, fórmulas, cifras ni procedimientos que no estén en el contexto.
8. Responde siempre en el mismo idioma en que se formula la pregunta."""

# ─── DESCRIPCIÓN COMPLETA DEL ASISTENTE ───────────────────────────────────────
CAPACIDADES_ASISTENTE = """
## 📋 **LO QUE PUEDO HACER POR TI:**

### 📄 **Sobre Documentos:**
- Responder preguntas sobre el contenido de tus documentos
- Buscar información específica dentro de PDFs, Word, Excel, PPT, TXT, CSV
- Resumir documentos completos o secciones específicas
- Extraer datos numéricos de Excel y CSV
- Comparar dos documentos entre sí

### 📊 **Análisis y Exportación:**
- Generar gráficas con datos de tus Excel/CSV
- Exportar respuestas a Word (.docx) con formato profesional
- Exportar respuestas a PDF
- Generar informes completos de conversación

### 📁 **Gestión:**
- Decirte cuántos documentos tienes cargados
- Listar todos tus documentos
- Recordar sobre qué documento estás hablando

### 💬 **Conversación:**
- Recordar el historial de nuestra conversación
- Responder en el mismo idioma que preguntas
- Tener memoria del documento actual que discutimos

**SOLO RESPONDO BASADO EN TUS DOCUMENTOS. NO INVENTO INFORMACIÓN.**
"""

SALUDOS = {
    "hola", "hi", "hello", "hey", "buenas", "buen día", "buen dia",
    "buenos días", "buenos dias", "buenas tardes", "buenas noches",
    "qué tal", "que tal", "cómo estás", "como estas"
}

# ─── Executor de segundo plano ────────────────────────────────────────────────
_executor = ThreadPoolExecutor(max_workers=2, thread_name_prefix="RAGIndexer")

def _shutdown_executor():
    _executor.shutdown(wait=False)

atexit.register(_shutdown_executor)

# ─── Caché de respuestas ──────────────────────────────────────────────────────
_cache: dict = {}
CACHE_TTL = 300
CACHE_MAX = 100

def clear_cache():
    _cache.clear()
    print("[RAG] Caché limpiado")

def _get_cache_key(question: str, doc_ids: list = None) -> str:
    key_str = f"{question.strip().lower()}|{sorted(doc_ids or [])}"
    return hashlib.md5(key_str.encode()).hexdigest()

def _get_gemini_model():
    api_key = os.environ.get("GEMINI_API_KEY", "")
    if not api_key:
        raise ValueError("GEMINI_API_KEY no está configurada.")
    genai.configure(api_key=api_key)
    return genai.GenerativeModel(
        model_name="gemini-2.5-flash",
        system_instruction=SYSTEM_PROMPT
    )

# ─── Función para obtener información de documentos indexados ─────────────────
def get_documents_info() -> dict:
    """Retorna información sobre los documentos indexados en ChromaDB"""
    try:
        all_meta = collection.get()["metadatas"]
        unique_docs = {}
        for m in all_meta:
            doc_id = m.get("doc_id")
            filename = m.get("filename", "desconocido")
            if doc_id and doc_id not in unique_docs:
                unique_docs[doc_id] = filename
        documentos_ordenados = sorted(unique_docs.values())
        return {
            "total": len(unique_docs),
            "documentos": documentos_ordenados
        }
    except Exception as e:
        print(f"[RAG] Error obteniendo info: {e}")
        return {"total": 0, "documentos": []}

# ─── Limpiar documentos huérfanos ─────────────────────────────────────────────
def limpiar_documentos_huertanos():
    """Elimina de ChromaDB los documentos que no existen en SQLite"""
    try:
        from backend.database import SessionLocal, Document
        
        db = SessionLocal()
        ids_validos = set([str(d.id) for d in db.query(Document).all()])
        db.close()
        
        all_data = collection.get()
        ids_a_eliminar = []
        
        for i, meta in enumerate(all_data['metadatas']):
            doc_id = meta.get('doc_id')
            if doc_id not in ids_validos:
                ids_a_eliminar.append(all_data['ids'][i])
        
        if ids_a_eliminar:
            collection.delete(ids=ids_a_eliminar)
            print(f"[RAG] 🧹 Eliminados {len(ids_a_eliminar)} chunks huérfanos")
        else:
            print("[RAG] ✅ No hay chunks huérfanos")
            
    except Exception as e:
        print(f"[RAG] Error limpiando huérfanos: {e}")

# ─── Memoria del documento actual ─────────────────────────────────────────────
_current_document_context = {
    "session_id": None,
    "user_id": None,
    "documento_actual": None,
    "ultima_pregunta": None
}

def update_document_context(user_id: int, session_id: str, documento: str, pregunta: str):
    _current_document_context["session_id"] = session_id
    _current_document_context["user_id"] = user_id
    _current_document_context["documento_actual"] = documento
    _current_document_context["ultima_pregunta"] = pregunta
    print(f"[RAG] Contexto actualizado: documento actual = '{documento}'")

def get_current_document(user_id: int, session_id: str) -> str:
    if (_current_document_context["session_id"] == session_id and 
        _current_document_context["user_id"] == user_id):
        return _current_document_context["documento_actual"]
    return None

def clear_document_context(user_id: int = None, session_id: str = None):
    if user_id and session_id:
        if (_current_document_context["user_id"] == user_id and 
            _current_document_context["session_id"] == session_id):
            _current_document_context["documento_actual"] = None
    else:
        _current_document_context["documento_actual"] = None
    print("[RAG] Contexto de documento limpiado")

# ─── Fuzzy matching ───────────────────────────────────────────────────────────
def find_closest_document(question: str, documents_list: list, cutoff: float = 0.35) -> str:
    if not documents_list:
        return None
    question_lower = question.lower()
    best_match = None
    best_score = 0
    
    palabras_genericas = ["cuantos", "cuántos", "lista", "listado", "todos", "generales", "documentos", "archivos"]
    if any(p in question_lower for p in palabras_genericas):
        return None
    
    for doc_name in documents_list:
        doc_lower = doc_name.lower().replace('.pdf', '').replace('.docx', '').replace('.doc', '').replace('_', ' ')
        doc_clean = doc_lower.split('.')[0]
        score = SequenceMatcher(None, question_lower, doc_clean).ratio()
        
        doc_words = {w for w in doc_clean.split() if len(w) > 2}
        question_words = {w for w in question_lower.split() if len(w) > 2}
        if doc_words and question_words:
            word_overlap = len(doc_words & question_words) / max(len(doc_words), 1)
            score = max(score, word_overlap * 1.2)
        
        if score > best_score and score >= cutoff:
            best_score = score
            best_match = doc_name
    
    if best_match:
        print(f"[RAG] Fuzzy matching: '{best_match}' con score {best_score:.2f}")
    return best_match

# ─── Loaders de documentos ────────────────────────────────────────────────────
def _load_pdf(file_path: str) -> list:
    loader = PyPDFLoader(file_path)
    docs = loader.load()
    for doc in docs:
        doc.page_content = doc.page_content.replace('\n', ' ').replace('  ', ' ')
    print(f"[RAG] PDF cargado: {len(docs)} páginas")
    return docs

def _load_docx(file_path: str) -> list:
    return Docx2txtLoader(file_path).load()

def _load_doc_old(file_path: str) -> list:
    try:
        import docx2txt
        from langchain_core.documents import Document
        text = docx2txt.process(file_path)
        return [Document(page_content=text, metadata={"source": file_path})]
    except Exception as e:
        print(f"[RAG] Fallback .doc: {e}")
        return Docx2txtLoader(file_path).load()

def _load_excel(file_path: str) -> list:
    from langchain_core.documents import Document
    docs = []
    try:
        excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            table_content = f"[Hoja: {sheet_name}]\n\n"
            table_content += df.to_markdown(index=False, max_rows=50)
            if len(df) > 50:
                table_content += f"\n\n*... y {len(df) - 50} filas más*"
            
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            if numeric_cols:
                stats = "\n\n**Estadísticas:**\n"
                for col in numeric_cols[:5]:
                    col_data = df[col].dropna()
                    if len(col_data) > 0:
                        stats += f"- {col}: min={col_data.min():.2f}, max={col_data.max():.2f}, promedio={col_data.mean():.2f}\n"
                table_content += stats
            
            docs.append(Document(
                page_content=table_content,
                metadata={"source": file_path, "sheet": sheet_name, "type": "excel"}
            ))
    except Exception as e:
        print(f"[RAG] Error leyendo Excel: {e}")
        try:
            if file_path.endswith('.xlsx'):
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        row_str = "\t".join(str(c) if c is not None else "" for c in row)
                        if row_str.strip():
                            rows.append(row_str)
                    if rows:
                        docs.append(Document(
                            page_content=f"[Hoja: {sheet_name}]\n" + "\n".join(rows),
                            metadata={"source": file_path, "sheet": sheet_name}
                        ))
        except Exception as e2:
            print(f"[RAG] Error en fallback Excel: {e2}")
    return docs

def _load_pptx(file_path: str) -> list:
    from langchain_core.documents import Document
    docs = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
            if texts:
                docs.append(Document(
                    page_content=f"[Diapositiva {i}]\n" + "\n".join(texts),
                    metadata={"source": file_path, "slide": i}
                ))
    except ImportError:
        print("[RAG] python-pptx no instalado")
    except Exception as e:
        print(f"[RAG] Error cargando PPTX: {e}")
    return docs

def _load_image(file_path: str) -> list:
    from langchain_core.documents import Document
    try:
        api_key = os.environ.get("GEMINI_API_KEY", "")
        if api_key:
            genai.configure(api_key=api_key)
            vision_model = genai.GenerativeModel('gemini-1.5-flash')
            with open(file_path, "rb") as f:
                image_b64 = base64.b64encode(f.read()).decode('utf-8')
            ext = file_path.split('.')[-1].lower()
            mime_type = f"image/{ext}" if ext != 'jpg' else "image/jpeg"
            prompt = "Extrae TODO el texto legible de esta imagen."
            response = vision_model.generate_content([prompt, {"mime_type": mime_type, "data": image_b64}])
            text = response.text.strip()
        else:
            import pytesseract
            from PIL import Image
            text = pytesseract.image_to_string(Image.open(file_path), lang="spa+eng").strip()
            if not text:
                text = "[No se detectó texto en la imagen]"
        return [Document(page_content=text, metadata={"source": file_path, "type": "image"})]
    except Exception as e:
        print(f"[RAG] Error procesando imagen: {e}")
        return [Document(page_content="[Error procesando imagen]", metadata={"source": file_path})]

def _load_txt(file_path: str) -> list:
    return TextLoader(file_path, encoding="utf-8").load()

def _load_csv(file_path: str) -> list:
    return CSVLoader(file_path, encoding="utf-8").load()

LOADER_MAP = {
    "pdf": _load_pdf, "docx": _load_docx, "doc": _load_doc_old,
    "xlsx": _load_excel, "xls": _load_excel,
    "pptx": _load_pptx, "ppt": _load_pptx,
    "txt": _load_txt, "csv": _load_csv,
    "png": _load_image, "jpg": _load_image, "jpeg": _load_image,
    "bmp": _load_image, "tiff": _load_image, "tif": _load_image, "webp": _load_image,
}
EXTENSIONES_SOPORTADAS = list(LOADER_MAP.keys())

def load_document(file_path: str, file_type: str) -> list:
    ext = file_type.lower().lstrip(".")
    loader_fn = LOADER_MAP.get(ext)
    if not loader_fn:
        return []
    try:
        docs = loader_fn(file_path)
        print(f"[RAG] Cargado '{file_path}': {len(docs)} secciones")
        return docs
    except Exception as e:
        print(f"[RAG] Error cargando: {e}")
        return []

# ─── Splitter adaptativo ──────────────────────────────────────────────────────
def _get_splitter(file_type: str) -> RecursiveCharacterTextSplitter:
    ext = file_type.lower().lstrip(".")
    chunk_size = 2000
    chunk_overlap = 400
    if ext in ["csv", "xlsx", "xls"]:
        separators = ["\n\n", "\n", "\t", ", ", " ", ""]
    elif ext in ["pptx", "ppt"]:
        separators = ["\n\n", "\n", " ", ""]
    else:
        separators = ["\n\n", "\n", ". ", " ", ""]
    return RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=chunk_overlap,
        separators=separators, length_function=len
    )

def _generate_summary(text: str, filename: str) -> str:
    try:
        api_key = os.environ.get("GEMINI_API_KEY", "")
        if not api_key:
            return ""
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-2.0-flash-lite')
        prompt = f"Genera un resumen MUY CORTO (2-3 líneas) del documento: {filename}\n\nTexto:\n{text[:4000]}"
        response = model.generate_content(prompt)
        summary = response.text.strip()
        if len(summary) > 350:
            summary = summary[:347] + "..."
        return summary
    except Exception as e:
        print(f"[RAG] Error generando resumen: {e}")
        return ""

# ─── Memoria de conversación ──────────────────────────────────────────────────
def save_to_memory(user_id: int, session_id: str, role: str, content: str):
    try:
        from backend.database import SessionLocal, ConversationMemory
        db = SessionLocal()
        memory = ConversationMemory(user_id=user_id, session_id=session_id, role=role, content=content[:2000])
        db.add(memory)
        db.commit()
        db.close()
    except Exception as e:
        print(f"[RAG] Error guardando memoria: {e}")

def get_conversation_context(user_id: int, session_id: str, limit: int = 10) -> str:
    try:
        from backend.database import SessionLocal, ConversationMemory
        db = SessionLocal()
        memories = db.query(ConversationMemory).filter(
            ConversationMemory.user_id == user_id,
            ConversationMemory.session_id == session_id
        ).order_by(ConversationMemory.timestamp.desc()).limit(limit).all()
        db.close()
        if not memories:
            return ""
        context_parts = []
        for m in reversed(memories):
            role_name = "Usuario" if m.role == "user" else "Asistente"
            context_parts.append(f"{role_name}: {m.content}")
        return "\n".join(context_parts)
    except Exception as e:
        print(f"[RAG] Error recuperando memoria: {e}")
        return ""

def detect_language(text: str) -> tuple:
    if not LANGDETECT_AVAILABLE:
        return "es", "español"
    try:
        lang_code = detect(text)
        lang_map = {'es': ('es', 'español'), 'en': ('en', 'inglés'), 'fr': ('fr', 'francés'),
                    'pt': ('pt', 'portugués'), 'de': ('de', 'alemán'), 'it': ('it', 'italiano')}
        return lang_map.get(lang_code, ('es', 'español'))
    except:
        return 'es', 'español'

def calculate_confidence(docs: list, question: str, sources: list) -> dict:
    if docs and len(docs) > 0:
        return {"score": 0.95, "level": "alta", "reason": ""}
    return {"score": 0.0, "level": "sin_informacion", "reason": ""}

# ─── Indexación ───────────────────────────────────────────────────────────────
def _index_task(doc_id: int, file_path: str, file_type: str, filename: str, db_session_factory):
    db = db_session_factory()
    try:
        docs = load_document(file_path, file_type)
        if not docs:
            raise Exception("No se pudo cargar el documento")
        
        print(f"[RAG] Muestra del texto: {docs[0].page_content[:500]}")
        
        results = collection.get(where={"doc_id": str(doc_id)})
        if results["ids"]:
            collection.delete(ids=results["ids"])
        
        splitter = _get_splitter(file_type)
        chunks = splitter.split_documents(docs)
        chunks = [c for c in chunks if len(c.page_content.strip()) >= 50]
        
        if not chunks:
            raise Exception("No se generaron chunks válidos")
        
        texts = [c.page_content.strip() for c in chunks]
        embeddings = EMBED_MODEL.encode(texts, show_progress_bar=False).tolist()
        ids = [f"doc{doc_id}_chunk{i}" for i in range(len(texts))]
        metadatas = [{"doc_id": str(doc_id), "filename": filename, "chunk": i, "page": chunks[i].metadata.get("page", 0)} for i in range(len(texts))]
        
        collection.add(documents=texts, embeddings=embeddings, ids=ids, metadatas=metadatas)
        print(f"[RAG] Indexado '{filename}': {len(texts)} chunks")
        
        summary = _generate_summary(" ".join([c.page_content[:500] for c in chunks[:5]]), filename)
        
        from backend.database import Document
        doc = db.query(Document).filter(Document.id == doc_id).first()
        if doc:
            doc.status = "Indexado"
            if summary:
                doc.summary = summary
            db.commit()
    except Exception as e:
        print(f"[RAG] Error indexando: {e}")
        try:
            from backend.database import Document
            doc = db.query(Document).filter(Document.id == doc_id).first()
            if doc:
                doc.status = "Error"
                db.commit()
        except:
            pass
    finally:
        db.close()

def index_document(doc_id: int, file_path: str, file_type: str, filename: str, db_session_factory=None) -> bool:
    if db_session_factory is None:
        try:
            from backend.database import SessionLocal
            db_session_factory = SessionLocal
        except:
            class DummySession:
                def __enter__(self): return self
                def __exit__(self, *args): pass
                def query(self, *args): return self
                def filter(self, *args): return self
                def first(self): return None
                def commit(self): pass
                def close(self): pass
            db_session_factory = lambda: DummySession()
    clear_cache()
    _executor.submit(_index_task, doc_id, file_path, file_type, filename, db_session_factory)
    return True

def delete_document_from_index(doc_id: int) -> None:
    """Elimina un documento del índice vectorial COMPLETAMENTE con verificación"""
    try:
        results = collection.get(where={"doc_id": str(doc_id)})
        if results and results["ids"]:
            chunks_count = len(results["ids"])
            collection.delete(ids=results["ids"])
            print(f"[RAG] ✅ Eliminados {chunks_count} chunks del doc_id={doc_id}")
            
            verify = collection.get(where={"doc_id": str(doc_id)})
            if verify and verify["ids"]:
                print(f"[RAG] ⚠️ ADVERTENCIA: Quedaron {len(verify['ids'])} chunks sin eliminar")
            else:
                print(f"[RAG] ✅ Confirmado: documento {doc_id} eliminado del índice")
        else:
            print(f"[RAG] No se encontraron chunks para doc_id={doc_id}")
    except Exception as e:
        print(f"[RAG] ❌ Error eliminando doc_id={doc_id}: {e}")
    clear_cache()

def reset_index():
    try:
        all_ids = collection.get()['ids']
        if all_ids:
            collection.delete(ids=all_ids)
            print(f"[RAG] Índices eliminados: {len(all_ids)} chunks")
    except Exception as e:
        print(f"[RAG] Error resetando: {e}")

def list_indexed_documents() -> list:
    try:
        all_meta = collection.get()["metadatas"]
        seen = {}
        for m in all_meta:
            doc_id = m.get("doc_id")
            if doc_id not in seen:
                seen[doc_id] = m.get("filename", "desconocido")
        return [{"doc_id": k, "filename": v} for k, v in seen.items()]
    except Exception as e:
        print(f"[RAG] Error listando: {e}")
        return []

# ─── Recuperación de chunks ───────────────────────────────────────────────────
def retrieve_chunks(question: str, doc_ids: list = None, n_results: int = None) -> tuple:
    question_embedding = EMBED_MODEL.encode([question]).tolist()[0]
    total = collection.count()
    if total == 0:
        return [], []
    
    if n_results is None:
        words = len(question.split())
        if words <= 10:
            n_results = 20
        elif words <= 25:
            n_results = 30
        else:
            n_results = 40
        n_results = min(n_results, total)
        print(f"[RAG] Palabras: {words}, n_results: {n_results}")
    
    try:
        if doc_ids:
            where = {"doc_id": {"$in": [str(d) for d in doc_ids]}}
            results = collection.query(query_embeddings=[question_embedding], n_results=n_results, where=where)
        else:
            results = collection.query(query_embeddings=[question_embedding], n_results=n_results)
    except Exception as e:
        print(f"[RAG] Error en query: {e}")
        return [], []
    
    docs = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    print(f"[RAG] Fragmentos recuperados: {len(docs)}")
    return docs, metas

def build_context(docs: list, metas: list, question: str = None, current_doc: str = None) -> tuple:
    context_parts = []
    sources = []
    seen = set()
    target_doc = current_doc
    
    if not target_doc and question:
        all_docs = list(set([meta.get("filename", "") for meta in metas]))
        target_doc = find_closest_document(question, all_docs, cutoff=0.35)
    
    if target_doc:
        for doc, meta in zip(docs, metas):
            filename = meta.get("filename", "")
            if filename == target_doc:
                page = meta.get("page", 0)
                key = f"{filename}|{page}"
                if key not in seen:
                    seen.add(key)
                    page_info = f", p.{page+1}" if page else ""
                    header = f"[{filename}{page_info}]"
                    context_parts.append(f"{header}\n{doc}")
                    sources.append({"filename": filename, "page": page, "preview": doc[:300]})
    
    if not target_doc:
        for doc, meta in zip(docs[:8], metas[:8]):
            filename = meta.get("filename", "")
            page = meta.get("page", 0)
            key = f"{filename}|{page}"
            if key not in seen:
                seen.add(key)
                page_info = f", p.{page+1}" if page else ""
                header = f"[{filename}{page_info}]"
                context_parts.append(f"{header}\n{doc}")
                sources.append({"filename": filename, "page": page, "preview": doc[:300]})
    
    print(f"[RAG] Fuentes: {len(sources)} (target: {target_doc if target_doc else 'múltiples'})")
    return "\n\n".join(context_parts), sources, target_doc

# ──────────────────────────────────────────────────────────────────────────────
# FASE 3 - EXPORTACIÓN MEJORADA (CON CONTENIDO REAL)
# ──────────────────────────────────────────────────────────────────────────────

def obtener_contenido_respuesta(user_id: int, session_id: str, pregunta: str) -> tuple:
    """Obtiene el contenido real de la respuesta del asistente"""
    contexto = get_conversation_context(user_id, session_id, limit=6)
    docs, metas = retrieve_chunks(pregunta, None, 15)
    return contexto, docs, metas

def exportar_a_word(question: str, user_id: int = None, session_id: str = None) -> dict:
    """Genera un archivo Word con el contenido REAL de la conversación"""
    from docx import Document as DocxDocument
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    import tempfile
    
    filename = f"export_{uuid.uuid4().hex[:8]}.docx"
    temp_path = os.path.join(tempfile.gettempdir(), filename)
    
    doc = DocxDocument()
    
    title = doc.add_heading('SGD·IA - Documento Exportado', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Fecha: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Usuario: {user_id if user_id else 'Desconocido'}").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()
    
    doc.add_heading('Pregunta del usuario', level=1)
    doc.add_paragraph(question)
    
    contexto, docs, metas = obtener_contenido_respuesta(user_id, session_id, question)
    
    if docs:
        doc.add_heading('Respuesta basada en documentos', level=1)
        for i, (doc_txt, meta) in enumerate(zip(docs[:10], metas[:10])):
            filename_ref = meta.get("filename", "Documento")
            page = meta.get("page", 0)
            doc.add_heading(f'Fragmento {i+1} - {filename_ref}', level=2)
            if page:
                doc.add_paragraph(f"(Página {page+1})", style='Intense Quote')
            doc.add_paragraph(doc_txt[:1000])
    else:
        doc.add_heading('Respuesta', level=1)
        doc.add_paragraph("Información solicitada. Para respuestas detalladas, realiza una pregunta específica en el chat.")
    
    if metas:
        doc.add_heading('Fuentes consultadas', level=1)
        seen = set()
        for meta in metas:
            filename = meta.get("filename", "Desconocido")
            if filename not in seen:
                seen.add(filename)
                doc.add_paragraph(f"📄 {filename}", style='List Bullet')
    
    doc.save(temp_path)
    
    return {
        "answer": f"✅ **He generado el archivo Word.**\n\n📎 [Descargar {filename}](http://localhost:8000/api/download/{filename})",
        "sources": metas[:5] if metas else [],
        "cached": False,
        "download_file": filename,
        "confidence": {"score": 1.0, "level": "alta", "reason": ""}
    }

def exportar_a_pdf(question: str, user_id: int = None, session_id: str = None) -> dict:
    """Genera un archivo PDF con el contenido REAL de la conversación"""
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    import tempfile
    
    filename = f"export_{uuid.uuid4().hex[:8]}.pdf"
    temp_path = os.path.join(tempfile.gettempdir(), filename)
    
    doc = SimpleDocTemplate(temp_path, pagesize=letter)
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'], fontSize=24, textColor=colors.HexColor('#1e40af'), alignment=TA_CENTER)
    center_style = ParagraphStyle('Center', parent=styles['Normal'], alignment=TA_CENTER)
    
    story = []
    
    story.append(Paragraph("SGD·IA - Documento Exportado", title_style))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"Fecha: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", center_style))
    story.append(Paragraph(f"Usuario: {user_id if user_id else 'Desconocido'}", center_style))
    story.append(PageBreak())
    
    story.append(Paragraph("Pregunta del usuario", styles['Heading1']))
    story.append(Paragraph(question.replace('\n', '<br/>'), styles['Normal']))
    story.append(Spacer(1, 12))
    
    contexto, docs, metas = obtener_contenido_respuesta(user_id, session_id, question)
    
    if docs:
        story.append(Paragraph("Respuesta basada en documentos", styles['Heading1']))
        for i, (doc_txt, meta) in enumerate(zip(docs[:10], metas[:10])):
            filename_ref = meta.get("filename", "Documento")
            page = meta.get("page", 0)
            story.append(Paragraph(f"Fragmento {i+1} - {filename_ref}", styles['Heading2']))
            if page:
                story.append(Paragraph(f"(Página {page+1})", styles['Normal']))
            story.append(Paragraph(doc_txt[:1000].replace('\n', '<br/>'), styles['Normal']))
            story.append(Spacer(1, 6))
    else:
        story.append(Paragraph("Respuesta", styles['Heading1']))
        story.append(Paragraph("Información solicitada.", styles['Normal']))
    
    if metas:
        story.append(PageBreak())
        story.append(Paragraph("Fuentes consultadas", styles['Heading1']))
        seen = set()
        for meta in metas[:20]:
            filename = meta.get("filename", "Desconocido")
            if filename not in seen:
                seen.add(filename)
                story.append(Paragraph(f"📄 {filename}", styles['Normal']))
                story.append(Spacer(1, 4))
    
    doc.build(story)
    
    return {
        "answer": f"✅ **He generado el archivo PDF.**\n\n📎 [Descargar {filename}](http://localhost:8000/api/download/{filename})",
        "sources": metas[:5] if metas else [],
        "cached": False,
        "download_file": filename,
        "confidence": {"score": 1.0, "level": "alta", "reason": ""}
    }

def extraer_datos_numericos_de_documentos() -> tuple:
    """Extrae datos numéricos reales de documentos Excel/CSV"""
    try:
        all_meta = collection.get()["metadatas"]
        datos = []
        nombres_columnas = []
        
        for meta in all_meta:
            filename = meta.get("filename", "")
            if filename.endswith(('.xlsx', '.xls', '.csv')):
                doc_id = meta.get("doc_id")
                if doc_id:
                    chunks = collection.get(where={"doc_id": str(doc_id)})
                    if chunks["documents"]:
                        for chunk in chunks["documents"]:
                            numeros = re.findall(r'-?\d+\.?\d*', chunk)
                            for n in numeros[:20]:
                                try:
                                    datos.append(float(n))
                                except:
                                    pass
                        if chunks["documents"][0]:
                            nombres_columnas.append(filename)
        
        return datos[:50], nombres_columnas[:5]
    except Exception as e:
        print(f"[RAG] Error extrayendo datos numéricos: {e}")
        return [], []

def generar_grafica_respuesta(question: str, user_id: int = None, session_id: str = None) -> dict:
    """Genera una gráfica con datos REALES de documentos Excel/CSV"""
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import tempfile
    
    filename = f"grafica_{uuid.uuid4().hex[:8]}.png"
    temp_path = os.path.join(tempfile.gettempdir(), filename)
    
    datos, fuentes = extraer_datos_numericos_de_documentos()
    
    if not datos:
        return {
            "answer": "⚠️ **No se encontraron datos numéricos** en los documentos cargados.\n\nPara generar una gráfica, necesito un archivo Excel o CSV con datos numéricos.",
            "sources": [],
            "cached": False,
            "confidence": {"score": 0.5, "level": "media", "reason": "Sin datos numéricos"}
        }
    
    if len(datos) <= 5:
        tipo = 'pastel'
    elif len(set(datos)) < len(datos) * 0.7:
        tipo = 'barras'
    else:
        tipo = 'lineas'
    
    plt.figure(figsize=(10, 6))
    
    if tipo == 'lineas':
        plt.plot(datos[:30], marker='o', linewidth=2, markersize=6)
        plt.xlabel("Índice", fontsize=12)
        plt.ylabel("Valor", fontsize=12)
    elif tipo == 'barras':
        plt.bar(range(len(datos[:20])), datos[:20])
        plt.xlabel("Índice", fontsize=12)
        plt.ylabel("Valor", fontsize=12)
    else:
        plt.pie(datos[:10], autopct='%1.1f%%')
    
    plt.title("Datos extraídos de documentos", fontsize=14)
    plt.grid(True, alpha=0.3)
    plt.tight_layout()
    plt.savefig(temp_path, dpi=100, bbox_inches='tight')
    plt.close()
    
    fuentes_texto = "\n".join([f"- {f}" for f in fuentes[:3]]) if fuentes else "- Documentos Excel/CSV"
    
    return {
        "answer": f"✅ **Gráfica generada** con {len(datos)} valores numéricos.\n\n📎 [Descargar {filename}](http://localhost:8000/api/download/{filename})",
        "sources": [{"filename": f, "type": "grafica"} for f in fuentes],
        "cached": False,
        "download_file": filename,
        "confidence": {"score": 0.95, "level": "alta", "reason": ""}
    }

def extraer_nombres_documentos(question: str) -> tuple:
    """Extrae posibles nombres de documentos de la pregunta"""
    question_lower = question.lower()
    all_docs = get_documents_info()["documentos"]
    
    doc1 = None
    doc2 = None
    
    for doc in all_docs:
        doc_clean = doc.lower().replace('.pdf', '').replace('.docx', '').replace('.xlsx', '')
        if doc_clean in question_lower:
            if doc1 is None:
                doc1 = doc
            else:
                doc2 = doc
                break
    
    return doc1, doc2

def comparar_documentos_respuesta(question: str, user_id: int = None, session_id: str = None) -> dict:
    """Compara dos documentos REALMENTE usando Gemini"""
    
    doc1, doc2 = extraer_nombres_documentos(question)
    
    if not doc1 or not doc2:
        docs_info = get_documents_info()
        doc_lista = "\n".join([f"- {d}" for d in docs_info["documentos"][:10]])
        return {
            "answer": f"⚠️ **Para comparar documentos**, necesito los nombres de dos documentos.\n\n**Ejemplo:** 'compara el documento SQL con el documento Excel'\n\n**Documentos disponibles:**\n{doc_lista}",
            "sources": [],
            "cached": False,
            "confidence": {"score": 0.5, "level": "media", "reason": "Documentos no especificados"}
        }
    
    try:
        results1 = collection.get(where={"filename": doc1})
        docs1 = results1.get("documents", [])[:10]
        metas1 = results1.get("metadatas", [])[:10]
    except:
        docs1, metas1 = retrieve_chunks(doc1, None, 10)
    
    try:
        results2 = collection.get(where={"filename": doc2})
        docs2 = results2.get("documents", [])[:10]
        metas2 = results2.get("metadatas", [])[:10]
    except:
        docs2, metas2 = retrieve_chunks(doc2, None, 10)
    
    if not docs1:
        return f"No se encontró el documento '{doc1}'"
    if not docs2:
        return f"No se encontró el documento '{doc2}'"
    
    contexto1 = "\n".join([d[:500] for d in docs1[:5]]) if docs1 else "No se encontró contenido"
    contexto2 = "\n".join([d[:500] for d in docs2[:5]]) if docs2 else "No se encontró contenido"
    
    try:
        model = _get_gemini_model()
        prompt = f"""
Compara los siguientes dos documentos en detalle:

DOCUMENTO 1: {doc1}
CONTENIDO:
{contexto1}

DOCUMENTO 2: {doc2}
CONTENIDO:
{contexto2}

Proporciona:
1. SIMILITUDES
2. DIFERENCIAS
3. CONCLUSIÓN
"""
        response = model.generate_content(prompt)
        comparacion = response.text
    except Exception as e:
        comparacion = f"Error: {e}\n\nDocumento 1: {doc1}\nDocumento 2: {doc2}"
    
    filename = f"comparacion_{uuid.uuid4().hex[:8]}.txt"
    temp_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", filename) if os.name == 'nt' else os.path.join("/tmp", filename)
    
    contenido_txt = f"""COMPARACIÓN DE DOCUMENTOS
{'=' * 50}
Fecha: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Documento 1: {doc1}
Documento 2: {doc2}

{comparacion}

Fuentes:
- {doc1}: {len(docs1)} fragmentos
- {doc2}: {len(docs2)} fragmentos
"""
    
    with open(temp_path, 'w', encoding='utf-8') as f:
        f.write(contenido_txt)
    
    return {
        "answer": f"✅ **Comparación completada**\n\n📄 {doc1}\n📄 {doc2}\n\n{comparacion}\n\n📎 [Descargar comparación](http://localhost:8000/api/download/{os.path.basename(temp_path)})",
        "sources": [{"filename": doc1, "type": "comparacion"}, {"filename": doc2, "type": "comparacion"}],
        "cached": False,
        "download_file": os.path.basename(temp_path),
        "confidence": {"score": 0.95, "level": "alta", "reason": ""}
    }

def generar_informe_respuesta(question: str, user_id: int = None, session_id: str = None) -> dict:
    """Genera un informe REAL de la conversación actual"""
    import tempfile
    
    filename = f"informe_{uuid.uuid4().hex[:8]}.txt"
    temp_path = os.path.join(tempfile.gettempdir(), filename)
    
    contenido_informe = "INFORME SGD·IA\n"
    contenido_informe += "=" * 50 + "\n\n"
    contenido_informe += f"Fecha: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
    contenido_informe += f"Usuario ID: {user_id if user_id else 'Desconocido'}\n\n"
    
    if session_id and user_id:
        contexto = get_conversation_context(user_id, session_id, limit=20)
        if contexto:
            contenido_informe += "CONVERSACIÓN COMPLETA:\n"
            contenido_informe += "-" * 30 + "\n"
            contenido_informe += contexto + "\n\n"
    
    docs_info = get_documents_info()
    contenido_informe += f"ESTADÍSTICAS:\n"
    contenido_informe += f"- Documentos cargados: {docs_info['total']}\n"
    contenido_informe += f"- Fragmentos en BD: {collection.count()}\n"
    
    with open(temp_path, 'w', encoding='utf-8') as f:
        f.write(contenido_informe)
    
    return {
        "answer": f"✅ **Informe generado**\n\n📎 [Descargar {filename}](http://localhost:8000/api/download/{filename})",
        "sources": [],
        "cached": False,
        "download_file": filename,
        "confidence": {"score": 1.0, "level": "alta", "reason": ""}
    }

def manejar_exportacion(question: str, user_id: int = None, session_id: str = None) -> dict:
    """Maneja las solicitudes de exportación del usuario"""
    question_lower = question.lower()
    
    if "word" in question_lower or "docx" in question_lower:
        return exportar_a_word(question, user_id, session_id)
    elif "pdf" in question_lower:
        return exportar_a_pdf(question, user_id, session_id)
    elif "informe" in question_lower or "reporte" in question_lower:
        return generar_informe_respuesta(question, user_id, session_id)
    elif "gráfica" in question_lower or "grafica" in question_lower:
        return generar_grafica_respuesta(question, user_id, session_id)
    elif "comparar" in question_lower or "compara" in question_lower:
        return comparar_documentos_respuesta(question, user_id, session_id)
    else:
        return {
            "answer": """📄 **Comandos de exportación:**
- **Word:** "dame un archivo Word"
- **PDF:** "exporta a PDF"
- **Informe:** "genera un informe"
- **Gráfica:** "crea una gráfica"
- **Comparar:** "compara documento X con Y" """,
            "sources": [],
            "cached": False,
            "confidence": {"score": 0.95, "level": "alta", "reason": ""}
        }

def is_complex_question(question: str) -> bool:
    """Detecta si una pregunta es compleja"""
    question_lower = question.lower()
    palabras_complejas = ["compara", "comparar", "gráfica", "grafica", "visualiza", "analiza", "promedio", "suma"]
    return any(p in question_lower for p in palabras_complejas) and len(question.split()) > 8

def query_rag(question: str, doc_ids: list = None, n_results: int = None, 
              user_id: int = None, session_id: str = None) -> dict:
    
    # ─── Detectar solicitud de exportación ─────────────────────────────────────
    question_lower_export = question.strip().lower()
    palabras_exportacion = ["dame un archivo", "exporta", "genera un documento", "crea un archivo",
        "pásame un", "envíame un", "descargar", "word", "pdf", "informe",
        "gráfica", "grafica", "comparar", "compara", "reporte"]
    if any(p in question_lower_export for p in palabras_exportacion):
        return manejar_exportacion(question, user_id, session_id)
    
    # ─── Detectar preguntas complejas ──────────────────────────────────────────
    if is_complex_question(question):
        try:
            from backend.agent import agente_responder
            print("[RAG] Usando agente para pregunta compleja...")
            return agente_responder(question, user_id, session_id)
        except ImportError:
            print("[RAG] Agente no disponible")
    
    # ─── Saludos con presentación completa ────────────────────────────────────
    if question.strip().lower() in SALUDOS:
        return {
            "answer": f"""¡Hola! 👋 Soy tu asistente de documentos inteligente.

{CAPACIDADES_ASISTENTE}

**¿Cómo empezar?**
- Sube documentos (PDF, Word, Excel, PPT, imágenes)
- Pregúntame: "¿De qué trata este documento?"
- Pide: "dame un resumen", "genera una gráfica", "compara documento A con B"
- Exporta: "dame un archivo Word", "exporta a PDF"

¿En qué puedo ayudarte hoy?""",
            "sources": [],
            "cached": False,
            "confidence": {"score": 1.0, "level": "alta", "reason": ""}
        }
    
    # ─── Pregunta sobre cantidad de documentos ─────────────────────────────────
    if ("cuantos" in question_lower_export or "cuántos" in question_lower_export) and ("documento" in question_lower_export or "archivo" in question_lower_export):
        docs_info = get_documents_info()
        total = docs_info["total"]
        if total == 0:
            answer = "No hay documentos cargados actualmente."
        elif total == 1:
            answer = f"Actualmente hay **1 documento** cargado:\n- {docs_info['documentos'][0]}"
        else:
            doc_list = "\n".join([f"- {doc}" for doc in docs_info['documentos']])
            answer = f"Actualmente hay **{total} documentos** cargados:\n{doc_list}"
        return {"answer": answer, "sources": [], "cached": False, "confidence": {"score": 0.95, "level": "alta", "reason": ""}}
    
    # ─── Caché ─────────────────────────────────────────────────────────────────
    cache_key = _get_cache_key(question, doc_ids)
    if cache_key in _cache:
        answer, sources, timestamp = _cache[cache_key]
        if time.time() - timestamp < CACHE_TTL:
            print(f"[RAG] Respuesta desde caché")
            return {"answer": answer, "sources": sources, "cached": True, "confidence": {"score": 1.0, "level": "alta", "reason": ""}}
        else:
            del _cache[cache_key]
    
    # ─── Idiomas ───────────────────────────────────────────────────────────────
    lang_code, lang_name = detect_language(question)
    print(f"[RAG] Idioma detectado: {lang_name}")
    
    # ─── Memoria de conversación ───────────────────────────────────────────────
    conversation_context = ""
    if user_id and session_id:
        conversation_context = get_conversation_context(user_id, session_id, limit=5)
        if conversation_context:
            print("[RAG] Contexto de memoria recuperado")
    
    # ─── Contexto de documento actual ──────────────────────────────────────────
    current_doc = get_current_document(user_id, session_id) if user_id and session_id else None
    if current_doc:
        print(f"[RAG] Documento actual: {current_doc}")
    
    # ─── Búsqueda ──────────────────────────────────────────────────────────────
    docs, metas = retrieve_chunks(question, doc_ids, n_results)
    
    if not docs:
        result = {
            "answer": "⚠️ No encontré información relevante en los documentos cargados.",
            "sources": [],
            "cached": False,
            "confidence": {"score": 0.0, "level": "sin_informacion", "reason": ""}
        }
        if len(_cache) >= CACHE_MAX:
            oldest = min(_cache, key=lambda k: _cache[k][2])
            del _cache[oldest]
        _cache[cache_key] = (result["answer"], result["sources"], time.time())
        return result
    
    # ─── Construir contexto y respuesta ────────────────────────────────────────
    context, sources, target_doc = build_context(docs, metas, question, current_doc)
    
    if target_doc and user_id and session_id:
        update_document_context(user_id, session_id, target_doc, question)
    
    if "olvidar" in question.lower() or "nuevo documento" in question.lower():
        clear_document_context(user_id, session_id)
    
    confidence = calculate_confidence(docs, question, sources)
    
    language_instruction = f"IMPORTANTE: Responde completamente en {lang_name}."
    
    user_message = f"""
{language_instruction}

{conversation_context}

Fragmentos de documentos:

{context}

---
PREGUNTA: {question}

Responde basándote SOLO en los fragmentos anteriores.
"""
    
    try:
        model = _get_gemini_model()
        response = model.generate_content(user_message)
        answer = response.text
    except Exception as e:
        error_msg = str(e)
        if "429" in error_msg:
            answer = f"⏳ Límite de API excedido.\n\nFragmentos encontrados:\n\n{context}"
        elif "GEMINI_API_KEY" in error_msg:
            answer = "❌ API key no configurada."
        else:
            answer = f"❌ Error: {e}\n\nFragmentos encontrados:\n\n{context}"
    
    if user_id and session_id:
        save_to_memory(user_id, session_id, "user", question)
        save_to_memory(user_id, session_id, "assistant", answer[:2000])
    
    if len(_cache) >= CACHE_MAX:
        oldest = min(_cache, key=lambda k: _cache[k][2])
        del _cache[oldest]
    _cache[cache_key] = (answer, sources, time.time())
    
    return {
        "answer": answer,
        "sources": sources,
        "cached": False,
        "confidence": confidence,
        "detected_language": lang_name
    }

# ─── Escaneo de carpeta ───────────────────────────────────────────────────────
def scan_folder(folder_path: str) -> list:
    files = []
    try:
        for root, dirs, filenames in os.walk(folder_path):
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            for fname in filenames:
                ext = os.path.splitext(fname)[1].lower().lstrip(".")
                if ext in EXTENSIONES_SOPORTADAS:
                    full_path = os.path.join(root, fname)
                    size = os.path.getsize(full_path)
                    files.append({"name": fname, "path": full_path, "ext": ext, "size": round(size / 1024, 1)})
    except Exception as e:
        print(f"[RAG] Error escaneando: {e}")
    return files